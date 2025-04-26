from pptx import Presentation

# Load the presentation
presentation_path = 'data/gaia/2023/validation/a3fbeb63-0e8c-4a11-bff6-0e3b484c3e9c.pptx'
presentation = Presentation(presentation_path)

# Extract text from each slide
slide_texts = []
for slide in presentation.slides:
    slide_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            slide_text.append(shape.text)
    slide_texts.append(" ".join(slide_text))

# Count slides mentioning crustaceans
crustacean_count = sum('crustacean' in text.lower() for text in slide_texts)
crustacean_count